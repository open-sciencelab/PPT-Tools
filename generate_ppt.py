import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pypinyin import lazy_pinyin
import tkinter as tk
from tkinter import filedialog
import os

def read_names_from_excel(excel_path):
    df = pd.read_excel(excel_path)
    return df['中文名'].tolist()

def generate_pinyin(name):
    return ''.join(lazy_pinyin(name)).upper()

def get_template_type(template_ppt):
    """
    根据模板文件名返回模板类型。
    """
    template_name = os.path.basename(template_ppt)

    if template_name == "eg1.pptx":
        return "type1"
    elif template_name == "eg2.pptx":
        return "type2"
    else:
        raise ValueError(f"未知模板文件: {template_ppt}")

def create_ppt_from_template(excel_file, template_ppt, output_ppt):
    prs = Presentation(template_ppt)
    template_slide = prs.slides[0]
    slide_layout = template_slide.slide_layout

    # 根据模板文件名获取模板类型
    template_type = get_template_type(template_ppt)

    names = read_names_from_excel(excel_file)

    for name in names:
        pinyin = generate_pinyin(name)
        new_slide = prs.slides.add_slide(slide_layout)

        # 根据模板类型设置占位符和字体样式
        if template_type == "type1":
            zh_placeholder_1_idx = 10
            zh_placeholder_2_idx = 11
            en_placeholder_1_idx = 13
            en_placeholder_2_idx = 14
            zh_font_size = Pt(73.8)
            zh_font_bold = True
            en_font_size = Pt(28)

        elif template_type == "type2":
            zh_placeholder_1_idx = 10
            zh_font_bold = True
            zh_font_size = Pt(72)
            # 对名字进行处理：如果是两个字，添加三个空格
            if len(name) == 2:
                name = f"{name[0]}   {name[1]}"

        # 中文占位符1
        zh_placeholder_1 = new_slide.placeholders[zh_placeholder_1_idx]
        zh_placeholder_1.text = name
        for run in zh_placeholder_1.text_frame.paragraphs[0].runs:
            run.font.size = zh_font_size
            run.font.bold = zh_font_bold

        # 如果模板是 type1，则设置第二个中文占位符和拼音占位符
        if template_type == "type1":
            # 中文占位符2
            zh_placeholder_2 = new_slide.placeholders[zh_placeholder_2_idx]
            zh_placeholder_2.text = name
            for run in zh_placeholder_2.text_frame.paragraphs[0].runs:
                run.font.size = zh_font_size
                run.font.bold = zh_font_bold

            # 拼音占位符1
            en_placeholder_1 = new_slide.placeholders[en_placeholder_1_idx]
            en_placeholder_1.text = pinyin
            for run in en_placeholder_1.text_frame.paragraphs[0].runs:
                run.font.size = en_font_size

            # 拼音占位符2
            en_placeholder_2 = new_slide.placeholders[en_placeholder_2_idx]
            en_placeholder_2.text = pinyin
            for run in en_placeholder_2.text_frame.paragraphs[0].runs:
                run.font.size = en_font_size

    prs.save(output_ppt)
    print(f"PPT已保存至: {output_ppt}")

def select_file(filetype, title):
    """
    打开文件对话框，选择文件，并返回文件路径
    :param filetype: 允许选择的文件类型
    :param title: 对话框标题
    :return: 选择的文件路径
    """
    file_path = filedialog.askopenfilename(filetypes=filetype, title=title)
    return file_path

def main():
    # 创建 Tkinter 根窗口
    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口

    # 用户选择模板
    template_ppt = filedialog.askopenfilename(title="选择模板文件", filetypes=[("PPTX files", "*.pptx")])
    excel_file = select_file([("Excel files", "*.xlsx")], "请选择Excel文件")
    output_ppt = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPTX files", "*.pptx")])

    create_ppt_from_template(excel_file, template_ppt, output_ppt)

if __name__ == "__main__":
    main()
