import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pypinyin import lazy_pinyin
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
import shutil
import traceback
import subprocess
import sys

from pptx.dml.color import RGBColor

def apply_text_style(run, font_name="微软雅黑", font_size=Pt(72), bold=True):
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)  # 白色字体

def resource_path(relative_path):
    """获取打包后资源的绝对路径"""
    try:
        base_path = sys._MEIPASS  # 打包后的临时目录
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

TEMPLATES = {
    "席卡": resource_path("eg1.pptx"),
    "姓名卡": resource_path("eg2.pptx")
}


def clean_name(name):
    name = str(name).strip().replace(" ", "").replace("\t", "")
    if len(name) == 2:
        name = f"{name[0]}  {name[1]}"  # 两个字中间加两个空格
    return name

# def read_names_from_excel(excel_path):
#     df = pd.read_excel(excel_path)
#     return [clean_name(name) for name in df['中文名'].dropna()]

def read_names_from_excel(excel_path):
    df = pd.read_excel(excel_path, header=None)  # 不把第一行当列名
    first_column = df.iloc[:, 0]  # 选取第1列（不管列名是什么）
    return [clean_name(name) for name in first_column.dropna()]


def parse_names_from_text(raw_text):
    names = raw_text.strip().splitlines()
    return [clean_name(name) for name in names if clean_name(name)]

def generate_pinyin(name):
    return ''.join(lazy_pinyin(name.replace(" ", ""))).upper()

def get_template_type(template_ppt):
    template_name = os.path.basename(template_ppt)
    if template_name == "eg1.pptx":
        return "type1"
    elif template_name == "eg2.pptx":
        return "type2"
    else:
        raise ValueError(f"未知模板文件: {template_ppt}")

# def build_presentation(names, template_ppt, include_pinyin):
#     prs = Presentation(template_ppt)
#     template_slide = prs.slides[0]
#     slide_layout = template_slide.slide_layout
#     template_type = get_template_type(template_ppt)

#     for name in names:
#         pinyin = generate_pinyin(name)
#         new_slide = prs.slides.add_slide(slide_layout)

#         if template_type == "type1":
#             zh_placeholder_1_idx = 10
#             zh_placeholder_2_idx = 11
#             en_placeholder_1_idx = 13
#             en_placeholder_2_idx = 14
#             zh_font_size = Pt(73.8)
#             zh_font_bold = True
#             en_font_size = Pt(28)
#         elif template_type == "type2":
#             zh_placeholder_1_idx = 10
#             zh_font_bold = True
#             zh_font_size = Pt(72)

#         zh_placeholder_1 = new_slide.placeholders[zh_placeholder_1_idx]
#         zh_placeholder_1.text = name
#         for run in zh_placeholder_1.text_frame.paragraphs[0].runs:
#             run.font.size = zh_font_size
#             run.font.bold = zh_font_bold

#         if template_type == "type1":
#             zh_placeholder_2 = new_slide.placeholders[zh_placeholder_2_idx]
#             zh_placeholder_2.text = name
#             for run in zh_placeholder_2.text_frame.paragraphs[0].runs:
#                 run.font.size = zh_font_size
#                 run.font.bold = zh_font_bold

#             if include_pinyin:
#                 en_placeholder_1 = new_slide.placeholders[en_placeholder_1_idx]
#                 en_placeholder_1.text = pinyin
#                 for run in en_placeholder_1.text_frame.paragraphs[0].runs:
#                     run.font.size = en_font_size

#                 en_placeholder_2 = new_slide.placeholders[en_placeholder_2_idx]
#                 en_placeholder_2.text = pinyin
#                 for run in en_placeholder_2.text_frame.paragraphs[0].runs:
#                     run.font.size = en_font_size

#     # ✅ 删除模板第一页，只在全部新幻灯片添加之后
#     prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
#     return prs

def build_presentation(names, template_ppt, include_pinyin):
    prs = Presentation(template_ppt)
    template_slide = prs.slides[0]
    slide_layout = template_slide.slide_layout
    template_type = get_template_type(template_ppt)

    if template_type == "type1":
        for name in names:
            pinyin = generate_pinyin(name)
            new_slide = prs.slides.add_slide(slide_layout)

            zh_placeholder_1_idx = 10
            zh_placeholder_2_idx = 11
            en_placeholder_1_idx = 13
            en_placeholder_2_idx = 14
            zh_font_size = Pt(73.8)
            zh_font_bold = True
            en_font_size = Pt(28)

            zh_placeholder_1 = new_slide.placeholders[zh_placeholder_1_idx]
            zh_placeholder_1.text = name
            for run in zh_placeholder_1.text_frame.paragraphs[0].runs:
                run.font.size = zh_font_size
                run.font.bold = zh_font_bold

            zh_placeholder_2 = new_slide.placeholders[zh_placeholder_2_idx]
            zh_placeholder_2.text = name
            for run in zh_placeholder_2.text_frame.paragraphs[0].runs:
                run.font.size = zh_font_size
                run.font.bold = zh_font_bold

            if include_pinyin:
                en_placeholder_1 = new_slide.placeholders[en_placeholder_1_idx]
                en_placeholder_1.text = pinyin
                for run in en_placeholder_1.text_frame.paragraphs[0].runs:
                    run.font.size = en_font_size

                en_placeholder_2 = new_slide.placeholders[en_placeholder_2_idx]
                en_placeholder_2.text = pinyin
                for run in en_placeholder_2.text_frame.paragraphs[0].runs:
                    run.font.size = en_font_size

    elif template_type == "type2":
        for i in range(0, len(names), 3):
            chunk = names[i:i+3]
            new_slide = prs.slides.add_slide(slide_layout)

            placeholders = [
                ph for ph in new_slide.placeholders
                if ph.placeholder_format.type == 7  # BODY 类型
            ]

            if len(placeholders) < len(chunk):
                print(f"⚠️ 当前幻灯片只有 {len(placeholders)} 个姓名占位符")
                continue

            for j, name in enumerate(chunk):
                ph = placeholders[j]
                ph.text = name
                for run in ph.text_frame.paragraphs[0].runs:
                    apply_text_style(run)


    else:
        raise ValueError("未知模板类型")

    # 删除模板第一页
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
    return prs

def inspect_placeholders(ppt_path):
    prs = Presentation(ppt_path)
    slide = prs.slides[0]
    print(f"\n检查幻灯片: {ppt_path}")
    print(f"共 {len(slide.placeholders)} 个占位符:\n")

    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        name = shape.name
        try:
            text = shape.text
        except:
            text = "(无文本)"
        ph_type = shape.placeholder_format.type
        print(f" - idx={idx:<2} | name={name:<30} | type={ph_type:<10} | text={text}")


def cleanup_temp_files():
    try:
        ppt_path = os.path.join(os.getcwd(), "preview_temp.pptx")
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
        print("临时文件已清理")
    except Exception as e:
        print("清理失败:", e)

# ====== GUI 主逻辑 ======
def run_gui():
    prs_obj = {"ppt": None}

    def select_excel():
        path = filedialog.askopenfilename(filetypes=[("Excel files", "*.xlsx")], title="选择Excel文件")
        if path:
            excel_path.set(path)
    
    def clear_excel():
        excel_path.set("")
        messagebox.showinfo("清除成功", "已清除 Excel 文件选择。")

    def generate_preview():
        try:
            selected_template_key = template_var.get()
            template_file = TEMPLATES.get(selected_template_key)
            include_pinyin_flag = include_pinyin.get()

            if not template_file or not os.path.exists(template_file):
                messagebox.showerror("错误", "请选择有效的PPT模板")
                return

            names = []
            if excel_path.get() and os.path.exists(excel_path.get()):
                names = read_names_from_excel(excel_path.get())
            else:
                raw_text = name_text.get("1.0", "end")
                names = parse_names_from_text(raw_text)

            if not names:
                messagebox.showerror("错误", "请提供有效的姓名列表")
                return

            prs = build_presentation(names, template_file, include_pinyin_flag)

            prs_obj["ppt"] = prs

            temp_ppt_path = os.path.join(os.getcwd(), "preview_temp.pptx")
            prs.save(temp_ppt_path)

            # 直接打开PPT文件供用户查看和保存
            subprocess.Popen(['start', '', temp_ppt_path], shell=True)

            messagebox.showinfo("预览已打开", "PPT已在PowerPoint中打开，请手动保存或另存为。")

        except Exception as e:
            messagebox.showerror("生成失败", str(e))
            traceback.print_exc()

    def save_ppt():
        if not prs_obj["ppt"]:
            messagebox.showerror("错误", "请先生成预览")
            return

        output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPTX files", "*.pptx")])
        if not output_path:
            return

        try:
            prs_obj["ppt"].save(output_path)
            messagebox.showinfo("保存成功", f"PPT 已保存到：\n{output_path}")
            cleanup_temp_files()
        except Exception as e:
            messagebox.showerror("保存失败", str(e))
            traceback.print_exc()

    def on_closing():
        cleanup_temp_files()
        root.destroy()

    # inspect_placeholders("eg1.pptx")
    # inspect_placeholders("eg2.pptx")

    root = tk.Tk()
    root.title("🎓 PPT 批量生成工具")
    root.geometry("700x500")
    root.resizable(False, False)
    root.option_add("*Font", "微软雅黑 11")
    root.protocol("WM_DELETE_WINDOW", on_closing)

    style = ttk.Style(root)
    style.configure("TButton", padding=6)
    style.configure("TLabel", padding=6)

    mainframe = ttk.Frame(root, padding="20")
    mainframe.pack(expand=True)

    template_var = tk.StringVar(value=list(TEMPLATES.keys())[0])
    excel_path = tk.StringVar()
    include_pinyin = tk.BooleanVar(value=True)

    ttk.Label(mainframe, text="选择PPT模板: ").grid(row=0, column=0, sticky="e")
    ttk.Combobox(mainframe, textvariable=template_var, values=list(TEMPLATES.keys()), state="readonly", width=30).grid(row=0, column=1, sticky="w")

    ttk.Label(mainframe, text="选择Excel文件: ").grid(row=1, column=0, sticky="e", pady=(10, 0))

    excel_button_frame = ttk.Frame(mainframe)
    excel_button_frame.grid(row=1, column=1, columnspan=2, sticky="w", pady=(10, 0))

    ttk.Button(excel_button_frame, text="浏览", command=select_excel).pack(side="left")
    ttk.Button(excel_button_frame, text="❌ 清除 Excel", command=clear_excel).pack(side="left", padx=(10, 0))


    ttk.Label(mainframe, textvariable=excel_path, foreground="gray", wraplength=500).grid(row=2, column=0, columnspan=2, sticky="w")

    ttk.Label(mainframe, text="或粘贴姓名列表: ").grid(row=3, column=0, sticky="ne", pady=(10, 0))
    name_text = tk.Text(mainframe, height=6, width=50)
    name_text.grid(row=3, column=1, sticky="w", pady=(10, 0))

    ttk.Checkbutton(mainframe, text="生成拼音 (英文)", variable=include_pinyin).grid(row=4, column=0, columnspan=2, sticky="w", pady=(10, 0))

    ttk.Button(mainframe, text="📂 生成并打开PPT", command=generate_preview).grid(row=5, column=0, columnspan=2, pady=(20, 5))
    ttk.Button(mainframe, text="💾 另存为 PPT", command=save_ppt).grid(row=6, column=0, columnspan=2)

    root.mainloop()

if __name__ == "__main__":
    run_gui()
